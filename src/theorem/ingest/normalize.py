import csv
import io
import json
from pathlib import Path

from .envelope import Anchor, Envelope, Media, Table
from .sniff import sniff


class IngestError(Exception):
    """Raised when data cannot be ingested due to format issues or missing extras."""

    pass


def normalize(data: bytes, filename: str) -> Envelope:
    """Convert raw bytes to a normalized Envelope.

    Dispatches on detected format from sniff(). Produces Envelope with:
    - body: decoded text or empty for tables/images
    - tables: extracted structured data (CSV, JSON arrays, JSONL)
    - images: extracted media files
    - meta: format metadata
    """
    fmt = sniff(data, filename)

    if fmt == "binary":
        raise IngestError("Binary data cannot be ingested")

    if fmt == "zip":
        raise IngestError("ZIP archives cannot be ingested")

    if fmt == "markdown" or fmt == "text":
        text = data.decode("utf-8")
        return Envelope(
            body=text,
            meta={"format": fmt, "filename": filename},
        )

    if fmt == "csv":
        text = data.decode("utf-8")
        reader = csv.DictReader(io.StringIO(text))
        rows = [row for row in reader]
        table_name = Path(filename).stem
        return Envelope(
            body="",
            tables=[Table(name=table_name, rows=rows, origin="")],
            meta={"format": "csv", "filename": filename},
        )

    if fmt == "json":
        text = data.decode("utf-8")
        obj = json.loads(text)

        if (
            isinstance(obj, list)
            and obj
            and all(isinstance(item, dict) for item in obj)
        ):
            if _is_homogeneous_list_of_dicts(obj):
                rows = _dictlist_to_rows(obj)
                table_name = Path(filename).stem
                return Envelope(
                    body="",
                    tables=[Table(name=table_name, rows=rows, origin="")],
                    meta={"format": "json", "filename": filename},
                )

        pretty = json.dumps(obj, indent=2)
        body = f"```json\n{pretty}\n```"
        return Envelope(
            body=body,
            meta={"format": "json", "filename": filename},
        )

    if fmt == "jsonl":
        text = data.decode("utf-8")
        lines = text.strip().split("\n")
        objects = [json.loads(line) for line in lines if line.strip()]

        if objects:
            all_keys = set()
            for obj in objects:
                if isinstance(obj, dict):
                    all_keys.update(obj.keys())

            rows = []
            for obj in objects:
                if isinstance(obj, dict):
                    row = {k: str(obj.get(k, "")) for k in all_keys}
                    rows.append(row)

            table_name = Path(filename).stem
            return Envelope(
                body="",
                tables=[Table(name=table_name, rows=rows, origin="")],
                meta={"format": "jsonl", "filename": filename},
            )

        return Envelope(
            body="",
            meta={"format": "jsonl", "filename": filename},
        )

    if fmt == "png" or fmt == "jpeg" or fmt == "webp" or fmt == "gif":
        media = Media(
            data=data,
            format=fmt,
            meta={"bytes": len(data)},
            origin="",
        )
        return Envelope(
            body="",
            images=[media],
            meta={"format": fmt, "filename": filename},
        )

    if fmt == "pdf":
        try:
            import pdfplumber
        except ImportError as e:
            raise IngestError("PDF support needs: pip install theorem[pdf]") from e
        env = Envelope(meta={"format": "pdf", "filename": filename})
        parts: list[str] = []
        offset = 0
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            env.meta["pages"] = len(pdf.pages)
            for i, page in enumerate(pdf.pages, start=1):
                env.anchors.append(Anchor(offset=offset, page=i))
                text = page.extract_text() or ""
                parts.append(text)
                offset += len(text) + 2
                for t_i, rows in enumerate(page.extract_tables()):
                    if not rows or not rows[0]:
                        continue
                    header = [str(h or f"col{j}") for j, h in enumerate(rows[0])]
                    dict_rows = [
                        {header[j]: str(c or "") for j, c in enumerate(r)}
                        for r in rows[1:]
                    ]
                    env.tables.append(
                        Table(name=f"p{i}t{t_i}", rows=dict_rows, origin=f"page {i}")
                    )
        env.body = "\n\n".join(parts)
        return env

    if fmt == "docx":
        return _normalize_docx(data, filename)

    if fmt == "xlsx":
        return _normalize_xlsx(data, filename)

    if fmt == "pptx":
        return _normalize_pptx(data, filename)

    raise IngestError(f"Unknown format: {fmt}")


_OFFICE_EXTRA_MESSAGE = "Office format support requires: pip install 'theorem[office]'"


def _normalize_docx(data: bytes, filename: str) -> Envelope:
    try:
        import docx
        from docx.oxml.ns import qn
        from docx.table import Table as DocxTable
        from docx.text.paragraph import Paragraph
    except ImportError as e:
        raise IngestError(_OFFICE_EXTRA_MESSAGE) from e

    doc = docx.Document(io.BytesIO(data))
    body_lines: list[str] = []
    tables: list[Table] = []
    table_idx = 0

    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            text = para.text
            if not text:
                continue
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                level_str = style_name.replace("Heading", "").strip()
                level = int(level_str) if level_str.isdigit() else 1
                body_lines.append(f"{'#' * level} {text}")
            else:
                body_lines.append(text)
        elif child.tag == qn("w:tbl"):
            table_idx += 1
            docx_table = DocxTable(child, doc)
            rows_data = docx_table.rows
            if not rows_data:
                continue
            header = [cell.text for cell in rows_data[0].cells]
            rows = []
            for row in rows_data[1:]:
                cells = [cell.text for cell in row.cells]
                rows.append(
                    {
                        header[j]: str(cells[j])
                        for j in range(min(len(header), len(cells)))
                    }
                )
            tables.append(Table(name=f"table{table_idx}", rows=rows, origin=""))

    images: list[Media] = []
    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            image_part = rel.target_part
            images.append(
                Media(
                    data=image_part.blob,
                    format=image_part.content_type.split("/")[-1],
                    meta={"bytes": len(image_part.blob)},
                    origin=rel_id,
                )
            )

    return Envelope(
        body="\n\n".join(body_lines),
        tables=tables,
        images=images,
        meta={"format": "docx", "filename": filename},
    )


def _normalize_xlsx(data: bytes, filename: str) -> Envelope:
    try:
        import openpyxl
    except ImportError as e:
        raise IngestError(_OFFICE_EXTRA_MESSAGE) from e

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    tables: list[Table] = []
    for ws in wb.worksheets:
        rows_iter = ws.iter_rows(values_only=True)
        try:
            header_row = next(rows_iter)
        except StopIteration:
            continue
        header = [
            str(h) if h is not None else f"col{i}" for i, h in enumerate(header_row)
        ]
        rows = []
        for row in rows_iter:
            rows.append(
                {
                    header[i]: ("" if v is None else str(v))
                    for i, v in enumerate(row)
                    if i < len(header)
                }
            )
        tables.append(Table(name=ws.title, rows=rows, origin=f"sheet {ws.title}"))

    return Envelope(
        body="",
        tables=tables,
        meta={"format": "xlsx", "filename": filename},
    )


def _normalize_pptx(data: bytes, filename: str) -> Envelope:
    try:
        import pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        raise IngestError(_OFFICE_EXTRA_MESSAGE) from e

    prs = pptx.Presentation(io.BytesIO(data))
    body_lines: list[str] = []
    tables: list[Table] = []
    images: list[Media] = []
    table_idx = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        body_lines.append(f"# Slide {slide_idx}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                body_lines.append(shape.text_frame.text)
            if shape.has_table:
                table_idx += 1
                shape_table = shape.table
                rows_data = list(shape_table.rows)
                header = [cell.text for cell in rows_data[0].cells]
                rows = []
                for row in rows_data[1:]:
                    cells = [cell.text for cell in row.cells]
                    rows.append(
                        {
                            header[j]: str(cells[j])
                            for j in range(min(len(header), len(cells)))
                        }
                    )
                tables.append(
                    Table(
                        name=f"table{table_idx}",
                        rows=rows,
                        origin=f"slide {slide_idx}",
                    )
                )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                images.append(
                    Media(
                        data=image.blob,
                        format=image.ext,
                        meta={"bytes": len(image.blob)},
                        origin=f"slide {slide_idx}",
                    )
                )
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text:
                body_lines.append(f"Notes: {notes_text}")

    return Envelope(
        body="\n\n".join(body_lines),
        tables=tables,
        images=images,
        meta={"format": "pptx", "filename": filename},
    )


def _is_homogeneous_list_of_dicts(obj: list) -> bool:
    """Check if a list is homogeneous: all dicts with scalar values and overlapping keys.

    Overlapping means keys share at least one common key or are identical.
    Non-scalar values (lists, dicts, etc.) return False to keep as body.
    """
    if not obj or not all(isinstance(item, dict) for item in obj):
        return False

    if len(obj) == 1:
        return all(_is_scalar(v) for v in obj[0].values())

    all_keys = set()
    for item in obj:
        if not all(_is_scalar(v) for v in item.values()):
            return False
        all_keys.update(item.keys())

    if not all_keys:
        return False

    for item in obj:
        item_keys = set(item.keys())
        if not item_keys.intersection(all_keys):
            return False

    return True


def _is_scalar(value) -> bool:
    """Check if a value is a scalar (str, int, float, bool, None)."""
    return isinstance(value, (str, int, float, bool, type(None)))


def _dictlist_to_rows(obj: list) -> list[dict[str, str]]:
    """Convert list of dicts to rows with string values, union keys, fill missing with empty string."""
    all_keys = set()
    for item in obj:
        if isinstance(item, dict):
            all_keys.update(item.keys())

    rows = []
    for item in obj:
        if isinstance(item, dict):
            row = {k: str(item.get(k, "")) for k in all_keys}
            rows.append(row)
    return rows
